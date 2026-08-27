"""Generate the five-slide summary deck as an editable .pptx.

    python make_deck.py [out.pptx]

Everything is native PowerPoint shapes and text - including the architecture
diagram - so the deck can be edited normally rather than being a wall of images.
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# Palette - matches the HTML deck.
INK = RGBColor(0x1B, 0x21, 0x1F)
SOFT = RGBColor(0x4B, 0x54, 0x4F)
MUTED = RGBColor(0x6E, 0x77, 0x72)
RULE = RGBColor(0xD3, 0xD8, 0xD2)
PAPER = RGBColor(0xF3, 0xF4, 0xF1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SIGNAL = RGBColor(0x2D, 0x5F, 0x7C)
SIGNAL_BG = RGBColor(0xE3, 0xEC, 0xF1)
FLAG = RGBColor(0xA9, 0x69, 0x1A)
FLAG_BG = RGBColor(0xF5, 0xEA, 0xDA)
SUNKEN = RGBColor(0xE8, 0xEA, 0xE5)

DISPLAY = "Archivo"       # falls back to the theme sans if not installed
BODY = "Source Serif 4"
MONO = "Consolas"

W, H = Inches(13.333), Inches(7.5)
M = Inches(0.62)          # page margin


# ----------------------------------------------------------------- helpers

def autofit(tf):
    """Ask PowerPoint to shrink text that overruns its box.

    python-pptx cannot measure rendered text, so rather than guess at line counts
    we let PowerPoint do the fitting when it opens the file.
    """
    bodyPr = tf._txBody.bodyPr
    for tag in ("a:normAutofit", "a:spAutoFit", "a:noAutofit"):
        for el in bodyPr.findall(qn(tag)):
            bodyPr.remove(el)
    bodyPr.append(bodyPr.makeelement(qn("a:normAutofit"), {}))
    return tf


def textbox(slide, x, y, w, h, *, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def para(tf, text, *, size=12, font=BODY, color=SOFT, bold=False, space_after=6,
         space_before=0, caps=False, spacing=1.25, align=PP_ALIGN.LEFT, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    p.line_spacing = spacing
    r = p.add_run()
    r.text = text.upper() if caps else text
    f = r.font
    f.size, f.name, f.color.rgb, f.bold = Pt(size), font, color, bold
    return p


def rich(tf, parts, *, size=12, font=BODY, space_after=6, spacing=1.25, first=False):
    """parts = [(text, color, bold), ...] within one paragraph."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_after = Pt(space_after)
    p.line_spacing = spacing
    for text, color, bold in parts:
        r = p.add_run()
        r.text = text
        r.font.size, r.font.name, r.font.color.rgb, r.font.bold = Pt(size), font, color, bold
    return p


def box(slide, x, y, w, h, *, fill=WHITE, line=None, width=Pt(0.75), radius=0.04,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            s.adjustments[0] = radius
        except (IndexError, ValueError):
            pass
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = width
    s.shadow.inherit = False
    return s


def blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = PAPER
    return s


def header(slide, num, title):
    tf = textbox(slide, M, Inches(0.5), Inches(1.0), Inches(0.35))
    para(tf, num, size=11, font=MONO, color=SIGNAL, bold=True, first=True, space_after=0)
    tf = textbox(slide, M + Inches(0.95), Inches(0.42), W - M * 2 - Inches(0.95), Inches(0.55))
    para(tf, title, size=26, font=DISPLAY, color=INK, bold=True, caps=True,
         first=True, space_after=0, spacing=1.0)
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, M, Inches(1.12),
                                    W - M, Inches(1.12))
    ln.line.color.rgb = INK
    ln.line.width = Pt(1.5)


def card(slide, x, y, w, h, tag, title, body_parts=None, *, accent=SIGNAL):
    """A titled card. body_parts is a list of paragraphs (str) or (text,color,bold) lists."""
    box(slide, x, y, w, h, fill=WHITE)
    pad = Inches(0.22)
    tf = autofit(textbox(slide, x + pad, y + pad, w - pad * 2, h - pad * 2))
    para(tf, tag, size=8.5, font=MONO, color=accent, caps=True, first=True, space_after=5)
    para(tf, title, size=13.5, font=DISPLAY, color=INK, bold=True, space_after=6, spacing=1.1)
    for b in body_parts or []:
        if isinstance(b, str):
            para(tf, b, size=10.5, color=SOFT, space_after=5, spacing=1.3)
        else:
            rich(tf, b, size=10.5, space_after=5, spacing=1.3)
    return tf


def metric_strip(slide, x, y, w, items, *, h=Inches(0.86)):
    n = len(items)
    gap = Inches(0.08)
    cw = int((w - gap * (n - 1)) / n)
    for i, (v, l) in enumerate(items):
        cx = x + i * (cw + gap)
        box(slide, cx, y, cw, h, fill=WHITE)
        tf = textbox(slide, cx + Inches(0.14), y + Inches(0.13), cw - Inches(0.28), h - Inches(0.2))
        para(tf, v, size=19, font=DISPLAY, color=FLAG, bold=True, first=True,
             space_after=2, spacing=1.0)
        para(tf, l, size=7.5, font=MONO, color=MUTED, caps=True, space_after=0, spacing=1.0)


def swap_rows(tf, rows):
    """Before/after comparison lines inside a card."""
    for label, value, hot in rows:
        rich(tf, [(f"{label}   ", FLAG if hot else MUTED, hot),
                  (value, FLAG if hot else MUTED, hot)],
             size=9.5, font=MONO, space_after=3, spacing=1.15)


# ------------------------------------------------------------------ slides

def slide_title(prs):
    s = blank(prs)
    tf = textbox(s, M, Inches(1.55), Inches(9.4), Inches(0.3))
    para(tf, "Build review · Customer operations", size=10.5, font=MONO,
         color=SIGNAL, caps=True, first=True, space_after=0)

    tf = textbox(s, M, Inches(2.05), Inches(9.6), Inches(2.0))
    para(tf, "Intelligent Customer\nSignal Detector", size=52, font=DISPLAY,
         color=INK, bold=True, first=True, space_after=0, spacing=0.95)

    tf = textbox(s, M, Inches(4.15), Inches(7.6), Inches(1.0))
    para(tf, "A prototype that reads support text, product usage, billing and "
             "satisfaction data together, and tells a retention team who to call "
             "this week — and why.", size=15, color=SOFT, first=True, space_after=0,
         spacing=1.35)

    ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, M, Inches(5.45), W - M, Inches(5.45))
    ln.line.color.rgb = RULE
    ln.line.width = Pt(1)

    for i, (v, l) in enumerate([("0.851", "AUC on real outcomes"),
                                ("76.7%", "Top-10 worklist precision"),
                                ("42", "Median days of warning")]):
        x = M + Inches(i * 3.1)
        tf = textbox(s, x, Inches(5.72), Inches(2.9), Inches(0.85))
        para(tf, v, size=32, font=DISPLAY, color=FLAG, bold=True, first=True,
             space_after=3, spacing=1.0)
        para(tf, l, size=8.5, font=MONO, color=MUTED, caps=True, space_after=0, spacing=1.0)


def slide_1(prs):
    s = blank(prs)
    header(s, "01", "Problem understanding and objective")

    tf = textbox(s, M, Inches(1.42), Inches(11.6), Inches(0.6))
    para(tf, "The team finds out too late. A complaint gets escalated or a cancellation "
             "email lands, and by then the conversation that could have changed the outcome "
             "has already passed.", size=14.5, color=SOFT, first=True, space_after=0, spacing=1.3)

    col = Inches(5.85)
    card(s, M, Inches(2.32), col, Inches(1.92),
         "What the brief describes",
         "The signals already exist — nobody sees them together",
         ["Chat logs sit with support. Failed payments sit with finance. CSAT sits in a "
          "survey tool. Each is reviewed by hand, on its own schedule, by someone looking "
          "at a different screen. No single view says “this account is drifting.”"])

    card(s, M, Inches(4.42), col, Inches(1.92),
         "Where we sharpened it",
         "The angry customer isn’t the hard case",
         ["Someone who writes “we’re considering cancelling” has told you "
          "already. The account that actually gets missed is the quiet one: usage halved, "
          "nobody logged a ticket, renewal in seven weeks. Silence reads as health, and it "
          "isn’t."])

    x2 = M + col + Inches(0.3)
    tf = card(s, x2, Inches(2.32), col, Inches(4.02),
              "Objective",
              "Turn scattered signals into one ranked list a CSM can work on Monday morning",
              ["Three things it has to do:"])
    for n, (t, rest) in enumerate([
            ("Correlate", " behaviour, billing, satisfaction and language into a single "
                          "score, rather than alerting on each in isolation."),
            ("Prioritise", " — a ranked worklist, not a dashboard of 700 rows to interpret."),
            ("Explain", " — every flag arrives with the evidence behind it and a specific "
                        "next step, so it can be acted on without a second investigation.")], 1):
        rich(tf, [(f"0{n}  ", SIGNAL, True), (t, INK, True), (rest, SOFT, False)],
             size=10.5, space_after=7, spacing=1.3)


def slide_2(prs):
    s = blank(prs)
    header(s, "02", "Solution architecture and design flow")

    y0 = Inches(1.45)
    lane_lbl = [("Sources", M), ("Two parallel layers", Inches(3.3)),
                ("Scoring", Inches(8.35)), ("Output", Inches(11.0))]  # narrowed below
    for text, x in lane_lbl:
        tf = textbox(s, x, y0, min(Inches(2.4), W - M - x), Inches(0.22))
        para(tf, text, size=8, font=MONO, color=MUTED, caps=True, first=True, space_after=0)

    # sources
    src = [("usage / logins", 1.78), ("billing records", 2.42), ("CSAT / NPS", 3.06)]
    for label, top in src:
        b = box(s, M, Inches(top), Inches(2.3), Inches(0.48), fill=None, line=RULE, width=Pt(1))
        tf = b.text_frame
        tf.word_wrap = True
        para(tf, label, size=10, font=MONO, color=SOFT, first=True, space_after=0,
             align=PP_ALIGN.CENTER)
        b.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    b = box(s, M, Inches(4.42), Inches(2.3), Inches(0.48), fill=SIGNAL_BG, line=SIGNAL, width=Pt(1.25))
    para(b.text_frame, "support text", size=10, font=MONO, color=SIGNAL, bold=True,
         first=True, space_after=0, align=PP_ALIGN.CENTER)
    b.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # arrows into layers
    for top in (2.02, 2.66, 3.30):
        c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, M + Inches(2.3), Inches(top),
                                   Inches(3.3), Inches(2.66))
        c.line.color.rgb = RULE
        c.line.width = Pt(1.25)
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, M + Inches(2.3), Inches(4.66),
                               Inches(3.3), Inches(4.66))
    c.line.color.rgb = SIGNAL
    c.line.width = Pt(1.5)

    # deterministic layer
    box(s, Inches(3.3), Inches(1.72), Inches(4.8), Inches(1.95), fill=WHITE, line=RULE, width=Pt(1.25))
    tf = textbox(s, Inches(3.52), Inches(1.92), Inches(4.4), Inches(1.6))
    para(tf, "Deterministic layer — pandas", size=12, font=DISPLAY, color=INK,
         bold=True, first=True, space_after=6)
    para(tf, "usage_decline · support_strain · billing_friction\n"
             "satisfaction_drop · engagement_breadth", size=9.5, font=MONO, color=SOFT,
         space_after=6, spacing=1.25)
    para(tf, "each 0–1 with a human-readable evidence\nstring. never guesses.",
         size=9, font=MONO, color=MUTED, space_after=0, spacing=1.2)

    # llm layer
    box(s, Inches(3.3), Inches(3.82), Inches(4.8), Inches(1.95), fill=SIGNAL_BG,
        line=SIGNAL, width=Pt(1.5))
    tf = textbox(s, Inches(3.52), Inches(4.02), Inches(4.4), Inches(1.6))
    para(tf, "LLM layer — one call per account", size=12, font=DISPLAY, color=SIGNAL,
         bold=True, first=True, space_after=6)
    para(tf, "text_sentiment · churn_language\n+ rationale, action, owner", size=9.5,
         font=MONO, color=SOFT, space_after=6, spacing=1.25)
    para(tf, "strict JSON schema. sees the signals above,\nnever the score.",
         size=9, font=MONO, color=MUTED, space_after=0, spacing=1.2)

    # arrows to score
    for y1, y2, col, wdt in ((2.7, 3.5, RULE, 1.25), (4.8, 3.9, SIGNAL, 1.5)):
        c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(8.1), Inches(y1),
                                   Inches(8.35), Inches(y2))
        c.line.color.rgb = col
        c.line.width = Pt(wdt)

    # scoring
    box(s, Inches(8.35), Inches(2.55), Inches(2.4), Inches(1.35), fill=WHITE, line=INK, width=Pt(1.5))
    tf = textbox(s, Inches(8.5), Inches(2.72), Inches(2.1), Inches(1.1))
    para(tf, "weighted score", size=12, font=DISPLAY, color=INK, bold=True, first=True,
         space_after=4, align=PP_ALIGN.CENTER)
    para(tf, "Σ(weight × severity)", size=9.5, font=MONO, color=SOFT,
         space_after=4, align=PP_ALIGN.CENTER)
    para(tf, "renormalised over\npresent signals", size=8.5, font=MONO, color=MUTED,
         space_after=0, align=PP_ALIGN.CENTER, spacing=1.2)

    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(9.55), Inches(3.9),
                               Inches(9.55), Inches(4.2))
    c.line.color.rgb = RULE
    c.line.width = Pt(1.25)
    box(s, Inches(8.35), Inches(4.2), Inches(2.4), Inches(0.72), fill=WHITE, line=RULE, width=Pt(1))
    tf = textbox(s, Inches(8.5), Inches(4.34), Inches(2.1), Inches(0.5))
    para(tf, "banding + priority", size=9.5, font=MONO, color=INK, first=True,
         space_after=2, align=PP_ALIGN.CENTER)
    para(tf, "absolute or by rank", size=8.5, font=MONO, color=MUTED, space_after=0,
         align=PP_ALIGN.CENTER)

    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10.75), Inches(3.2),
                               Inches(11.0), Inches(3.2))
    c.line.color.rgb = RULE
    c.line.width = Pt(1.25)

    for i, label in enumerate(["watchlist", "heatmap", "account detail", "CSV export"]):
        b = box(s, Inches(11.0), Inches(2.2 + i * 0.58), Inches(1.72), Inches(0.46),
                fill=None, line=RULE, width=Pt(1))
        para(b.text_frame, label, size=9.5, font=MONO, color=SOFT, first=True,
             space_after=0, align=PP_ALIGN.CENTER)
        b.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # caption + the design decision
    tf = textbox(s, M, Inches(5.95), Inches(3.5), Inches(1.05))
    para(tf, "The two layers run in parallel and meet at the score. The model reads the "
             "computed signals alongside the raw text — but never sees or produces the "
             "risk score itself.",
         size=9, color=MUTED, first=True, space_after=0, spacing=1.28)

    card(s, Inches(4.35), Inches(5.78), Inches(4.1), Inches(1.32),
         "The decision this rests on",
         "The score comes from extracted evidence, not a number the model invented",
         ["Handing everything to the LLM would have worked and been unauditable — no way to "
          "answer “why is this account a 74?” beyond re-asking it."])

    card(s, Inches(8.62), Inches(5.78), Inches(4.1), Inches(1.32),
         "Why it ports",
         "A new source is a mapping function, not a rewrite",
         ["The detector’s input contract is two tables. Three differently shaped sources — "
          "wide, long, and a multi-table panel — ran through the same scoring code unchanged."])


def slide_3(prs):
    s = blank(prs)
    header(s, "03", "Implementation highlights")

    col = Inches(5.85)
    card(s, M, Inches(1.42), col, Inches(1.42),
         "AI logic", "One call per account, doing four jobs at once",
         ["Rather than separate calls for sentiment, intent, rationale and action, each "
          "account gets a single request returning all four under a strict JSON schema. "
          "Cheaper, faster, and the four outputs stay consistent."])

    card(s, M, Inches(3.02), col, Inches(1.72),
         "Prompt design", "Three rules doing the heavy lifting",
         ["Frustration is not exit intent — keeps escalations from reading as "
          "cancellations. Silence before a renewal is a risk, not health — inverts the "
          "model’s default read of a quiet account. Judge only from evidence given — "
          "why rationales quote real numbers instead of inventing history."])

    card(s, M, Inches(4.92), col, Inches(1.72),
         "Deliberate constraint", "No keyword fallback",
         ["An earlier version degraded to a lexicon when no API key was present. We removed "
          "it: a keyword matcher is not sentiment analysis. A missing key is now a hard stop; "
          "a failed call marks that account NOT ANALYSED rather than guessing a number."])

    x2 = M + col + Inches(0.3)
    tf = textbox(s, x2, Inches(1.42), col, Inches(0.25))
    para(tf, "Scoring, in full", size=8.5, font=MONO, color=MUTED, caps=True,
         first=True, space_after=0)

    box(s, x2, Inches(1.72), col, Inches(2.28), fill=WHITE)
    tf = textbox(s, x2 + Inches(0.22), Inches(1.9), col - Inches(0.44), Inches(2.0))
    code = [
        ("# seven signals, each normalised 0-1 with evidence", MUTED, False),
        ("WEIGHTS = {churn_language: .20, text_sentiment: .16,", INK, False),
        ("           usage_decline: .16, support_strain: .15,", INK, False),
        ("           billing_friction: .13, satisfaction_drop: .12,", INK, False),
        ("           engagement_breadth: .08}", INK, False),
        ("", INK, False),
        ("# renormalise over what this dataset carries -", MUTED, False),
        ("# a missing signal is skipped, never scored as zero", MUTED, False),
        ("total_w = sum(s.weight for s in sigs.values())", INK, False),
        ("raw = 100 * sum(s.weight * s.score ...) / total_w", INK, False),
        ("", INK, False),
        ("# renewal proximity is urgency, not risk", MUTED, False),
        ("amp = 1.25 if renewal_in <= 45 else 1.12 if <= 90 else 1.0", INK, False),
    ]
    for i, (line, colr, bold) in enumerate(code):
        para(tf, line or " ", size=8.5, font=MONO, color=colr, bold=bold,
             first=(i == 0), space_after=0, spacing=1.35)

    card(s, x2, Inches(4.18), col, Inches(1.42),
         "Validation harness", "A point-in-time backtest, not a scoreboard",
         ["backtest.py rewinds to a past date, shows the detector only records timestamped "
          "on or before it, then grades against what happened in the next 90 days. Three "
          "quarterly dates so horizons don’t overlap."])

    metric_strip(s, x2, Inches(5.78), col,
                 [("7", "signals fused"), ("2", "input formats"),
                  ("6", "self-checks"), ("~20s", "per 30 accounts")])


def slide_4(prs):
    s = blank(prs)
    header(s, "04", "Challenges and learnings")

    col = Inches(5.85)
    tf = card(s, M, Inches(1.42), col, Inches(2.32),
              "Challenge — calibration",
              "Fixed thresholds don’t survive a change of dataset",
              ["How severe a 20% usage drop looks depends entirely on the book. Absolute cut-"
               "points tuned on one population leave the top bands empty on a calmer one, even "
               "when the ranking is sound. So the detector can band by rank within the book "
               "instead — the worklist stays the right size whatever the data looks like."])
    swap_rows(tf, [("fixed thresholds", "bands empty · 12% recall", False),
                   ("banded by rank", "52% precision · 64% recall", True)])

    card(s, M, Inches(3.92), col, Inches(1.92),
         "Challenge — integration", "Field names promise more than they deliver",
         ["Every source we connected had at least one column whose name implied something "
          "different from what it held — a flag that meant “an event occurred” rather than "
          "“the customer left”, or an unknown outcome stored as a zero. Any production "
          "connection needs a validation pass before scoring, not after."])

    card(s, M, Inches(6.02), col, Inches(1.0),
         "Design principle", "Built for partial data, because production always is",
         ["No two feeds carry the same columns. Absent signals are skipped and their weight "
          "redistributed — never defaulted to zero, which reads as “healthy on that axis”."])

    x2 = M + col + Inches(0.3)
    card(s, x2, Inches(1.42), col, Inches(1.92),
         "Learning — measurement", "A terrible-looking metric that was a threshold bug",
         ["Sentiment classification scored 32% accuracy against human labels. Alarming, until "
          "we measured ranking separately: AUC 0.965. The model orders messages almost "
          "perfectly but compresses scores into 0.1–0.7, so fixed cut-points dumped most "
          "negatives into “neutral”."])

    tf = card(s, x2, Inches(3.52), col, Inches(2.52),
              "Learning — architecture", "The hybrid beats either approach on its own",
              ["We ran the ablation to find the right division of labour. Behavioural rules "
               "catch decay in the numbers but are blind to stated intent; the language layer "
               "reads intent but cannot see a usage cliff. Each misses what the other catches, "
               "and correlating them is what produces the result."])
    swap_rows(tf, [("rules only          ", "0.775 AUC", False),
                   ("LLM text only       ", "0.685 AUC", False),
                   ("hybrid, correlated  ", "0.851 AUC", True)])

    card(s, x2, Inches(6.22), col, Inches(0.8),
         "Net", "The hybrid is the product — +0.076 AUC over the best single layer, "
                "+0.196 over CSAT alone", [])


def slide_5(prs):
    s = blank(prs)
    header(s, "05", "Demo summary and next steps")

    col = Inches(5.85)
    card(s, M, Inches(1.42), col, Inches(1.55),
         "What’s running", "A Streamlit app over a seven-signal detector",
         ["Three tabs: a prioritised watchlist with per-account rationale and a recommended "
          "next step with an owner; a signal heatmap showing which lever to pull; and an "
          "account detail view breaking the score into points-of-risk beside the raw "
          "messages it was built from. Everything exports to CSV."])

    box(s, M, Inches(3.12), col, Inches(1.55), fill=WHITE)
    tf = textbox(s, M + Inches(0.22), Inches(3.3), col - Inches(0.44), Inches(1.3))
    cmds = [("# set OPENAI_API_KEY in .env, then", MUTED),
            ("streamlit run app.py        # the demo", INK),
            ("python signals.py           # CLI + evaluation", INK),
            ("python backtest.py          # outcome backtest", INK),
            ("python eval_llm.py          # LLM layer alone", INK),
            ("python test_signals.py      # self-checks, no key", INK)]
    for i, (line, colr) in enumerate(cmds):
        para(tf, line, size=9, font=MONO, color=colr, first=(i == 0), space_after=0, spacing=1.4)

    metric_strip(s, M, Inches(4.85), col,
                 [("0.851", "AUC, outcomes"), ("6.3x", "lift"),
                  ("42d", "lead time"), ("100%", "intent precision")])

    card(s, M, Inches(5.95), col, Inches(1.1),
         "Scope", "What a POC can and can’t settle",
         ["The backtest shows the method and the engineering hold up end to end. Real "
          "performance depends on the book it runs against, so the first step in any "
          "deployment is re-running this same harness on live history."])

    x2 = M + col + Inches(0.3)
    tf = card(s, x2, Inches(1.42), col, Inches(5.63),
              "Given more time — in the order we’d do it",
              "Close the gaps we already know about", [])
    steps = [
        ("Fold in intervention history.", " Knowing which accounts a CSM already reached "
         "lets the model tell a genuine miss from an account that was flagged and "
         "successfully saved — richer inputs sharpen precision."),
        ("Report uncertainty.", " Headline figures are point estimates. PR-AUC and bootstrap "
         "confidence intervals belong beside them, sized to whatever book it runs on."),
        ("Check stability over time.", " Results should hold at each scoring cycle "
         "separately, not just pooled across them."),
        ("Calibrate.", " Nothing yet says a score of 60 means a 60% chance of leaving. "
         "Ranking is validated; probability isn’t."),
        ("Fit the weights.", " They’re hand-set. With real churn outcomes they should be "
         "learned — the extraction layer wouldn’t change."),
        ("Run it continuously.", " Alert on band transitions rather than absolute scores, so "
         "the team hears about accounts that are moving."),
    ]
    for n, (b, rest) in enumerate(steps, 1):
        rich(tf, [(f"0{n}  ", SIGNAL, True), (b, INK, True), (rest, SOFT, False)],
             size=10.5, space_after=8, spacing=1.3)


def main(out: Path) -> None:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    for fn in (slide_title, slide_1, slide_2, slide_3, slide_4, slide_5):
        fn(prs)
    prs.save(out)
    print(f"{len(prs.slides.__iter__.__self__._sldIdLst)} slides -> {out}")


if __name__ == "__main__":
    main(Path(sys.argv[1]) if len(sys.argv) > 1 else
         Path(__file__).parent / "Signal_Detector_Summary.pptx")
